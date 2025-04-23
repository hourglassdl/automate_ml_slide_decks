from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

BASE_DIR = Path(__file__).resolve().parent
PPTX_PATH = BASE_DIR / "output/monthly_kpis.pptx"
IMG1_PATH = BASE_DIR / "output/at_risk_customers.png"
IMG2_PATH = BASE_DIR / "output/customer_breakdown_pie.png"
IMG3_PATH = BASE_DIR / "output/tenure_churn_predictions.png"

# Load PowerPoint
prs = Presentation()

# Choose a slide layout
slide_layout = prs.slide_layouts[6]  # Blank slide

# Add a new slide
slide = prs.slides.add_slide(slide_layout)

# Add title
left = top = Inches(0.5)
width = Inches(9)
height = Inches(1)

title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
title_frame.text = "Monthly KPIs"
title_frame.paragraphs[0].font.size = Pt(32)

# Add first image (at_risk_customers)
if IMG1_PATH.exists():
    slide.shapes.add_picture(str(IMG1_PATH), Inches(1), Inches(1.5), height=Inches(3))

# Add second image (customer_breakdown_pie)
if IMG2_PATH.exists():
    slide.shapes.add_picture(str(IMG2_PATH), Inches(5.5), Inches(1.5), height=Inches(3))

# Add third image (tenure_churn_predictions)
if IMG3_PATH.exists():
    slide.shapes.add_picture(str(IMG3_PATH), Inches(2), Inches(4), height=Inches(3.3))
    
# Save new presentation
prs.save(PPTX_PATH)
print(f"Presentation saved to {PPTX_PATH}")